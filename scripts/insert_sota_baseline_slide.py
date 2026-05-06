from pathlib import Path

from pptx import Presentation


DECK_PATH = Path("NV Compact Magnetometer IP Plan.pptx")
SLIDE_IMAGE = Path("output/slides/sota_baseline_slide.png")


def clear_slide(slide):
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)


def main():
    if not SLIDE_IMAGE.exists():
        raise FileNotFoundError(f"Missing rendered slide image: {SLIDE_IMAGE}")

    prs = Presentation(DECK_PATH)
    if len(prs.slides) < 2:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_id_list = prs.slides._sldIdLst
        created = slide_id_list[-1]
        slide_id_list.remove(created)
        slide_id_list.insert(1, created)
        slide = prs.slides[1]
    else:
        slide = prs.slides[1]
        clear_slide(slide)

    slide.shapes.add_picture(str(SLIDE_IMAGE), 0, 0, prs.slide_width, prs.slide_height)
    temp_path = DECK_PATH.with_suffix(".tmp.pptx")
    prs.save(temp_path)
    temp_path.replace(DECK_PATH)


if __name__ == "__main__":
    main()
